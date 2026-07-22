from __future__ import annotations

from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path
import math
import re

from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, PieChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


DOCS = Path.home() / "Documents"
SOURCE_BOARD = DOCS / "BOARD REPORT.xlsx"
SOURCE_SUMMARY = DOCS / "EXECUTIVE SUMMARY.docx"
SOURCE_TRACKER = DOCS / "JUNE 2026 CONSOLIDATED TRACKER.xlsx"

OUT_XLSX = DOCS / "ARCC Board Compliance Pack - Q2 2026 - Green Board Report v5.xlsx"
OUT_PPTX = DOCS / "ARCC Board Compliance Presentation - Q2 2026 - Green Board Report v5.pptx"
OUT_PDF = DOCS / "ARCC Board Compliance Presentation - Q2 2026 - Green Board Report v5.pdf"

COLORS = {
    "navy": RGBColor(13, 77, 45),
    "blue": RGBColor(28, 122, 72),
    "teal": RGBColor(57, 151, 94),
    "green": RGBColor(0, 128, 64),
    "amber": RGBColor(214, 141, 30),
    "red": RGBColor(176, 48, 48),
    "grey": RGBColor(91, 105, 125),
    "light": RGBColor(245, 247, 250),
    "white": RGBColor(255, 255, 255),
}


def clean_text(value) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\n", " ").replace("\r", " ")
    return re.sub(r"\s+", " ", text).strip()


def pct(value: float, digits: int = 0) -> str:
    return f"{value * 100:.{digits}f}%"


def status_key(value) -> str:
    raw = clean_text(value).lower().replace(".", "")
    if raw == "compliant":
        return "Compliant"
    if raw == "non compliant":
        return "Non Compliant"
    if raw in {"wip", "w i p"}:
        return "W.I.P"
    if raw == "n/a":
        return "N/A"
    return clean_text(value) or "Blank"


def safe_float(value):
    if value is None or value == "":
        return None
    try:
        return float(value)
    except Exception:
        return None


def format_metric(value) -> str:
    numeric = safe_float(value)
    if numeric is None:
        return clean_text(value) or "-"
    if 0 <= numeric <= 1:
        return pct(numeric, 0)
    if numeric == int(numeric):
        return f"{int(numeric):,}"
    return f"{numeric:,.2f}"


def rag_for_kpi(target, actual, status_note="") -> tuple[str, str]:
    note = clean_text(status_note)
    if "🔴" in note or note.lower() == "red":
        return "Red", "Target missed / board monitoring required"
    target_num = safe_float(target)
    actual_num = safe_float(actual)
    if target_num is None or actual_num is None:
        return ("Amber", note or "Monitor") if note else ("Green", "Completed")
    if target_num == 0:
        return ("Green", "No outstanding item") if actual_num == 0 else ("Red", f"{format_metric(actual)} outstanding item(s)")
    achievement = actual_num / target_num if target_num else 0
    if achievement >= 1:
        return "Green", "Target met"
    if achievement >= 0.8:
        return "Amber", "Close to target / monitor to closure"
    return "Red", "Below target"


def extract_summary_text() -> list[str]:
    doc = Document(SOURCE_SUMMARY)
    paragraphs = []
    for para in doc.paragraphs:
        text = clean_text(para.text)
        if text:
            paragraphs.append(text)
    return paragraphs


def load_board_dashboard():
    wb = load_workbook(SOURCE_BOARD, data_only=True)
    ws = wb["Sheet1"]
    dashboard = []
    for row in ws.iter_rows(min_row=10, max_row=20, values_only=True):
        area = clean_text(row[1])
        if not area:
            continue
        dashboard.append(
            {
                "sn": clean_text(row[0]),
                "area": area,
                "current": row[2],
                "previous": row[3],
                "comment": clean_text(row[5]),
                "board_attention": clean_text(row[6]),
            }
        )

    kpis = []
    ws2 = wb["Sheet2"]
    for row in ws2.iter_rows(min_row=4, values_only=True):
        if not clean_text(row[1]):
            continue
        kpis.append(
            {
                "kpi": clean_text(row[1]),
                "target": row[2],
                "actual": row[3],
                "status": clean_text(row[4]),
            }
        )
    return dashboard, kpis


def load_tracker():
    wb = load_workbook(SOURCE_TRACKER, data_only=True)
    ws = wb.active
    rows = []
    for row in ws.iter_rows(min_row=9, values_only=True):
        sn = row[1]
        item = clean_text(row[2])
        if sn is None or not item:
            continue
        rows.append(
            {
                "sn": sn,
                "item": item,
                "responsibility": clean_text(row[3]) or "Unassigned",
                "recipient": clean_text(row[4]),
                "frequency": clean_text(row[5]),
                "timeline": clean_text(row[6]),
                "remarks": clean_text(row[7]),
                "may_status": status_key(row[8]),
                "june_status": status_key(row[9]),
            }
        )
    return rows


def derive_pack_data():
    summary = extract_summary_text()
    dashboard, kpis = load_board_dashboard()
    tracker = load_tracker()
    counts = Counter(row["june_status"] for row in tracker)
    may_counts = Counter(row["may_status"] for row in tracker)
    by_owner = defaultdict(Counter)
    for row in tracker:
        by_owner[row["responsibility"]][row["june_status"]] += 1

    issues = [row for row in tracker if row["june_status"] in {"Non Compliant", "W.I.P"}]
    non_compliant = [row for row in tracker if row["june_status"] == "Non Compliant"]
    wip = [row for row in tracker if row["june_status"] == "W.I.P"]
    total = len(tracker)

    owner_rows = []
    for owner, counter in by_owner.items():
        open_count = counter["Non Compliant"] + counter["W.I.P"]
        if open_count:
            owner_rows.append(
                {
                    "owner": owner,
                    "compliant": counter["Compliant"],
                    "non_compliant": counter["Non Compliant"],
                    "wip": counter["W.I.P"],
                    "na": counter["N/A"],
                    "open": open_count,
                }
            )
    owner_rows.sort(key=lambda x: (-x["open"], x["owner"]))

    board_attention = [row for row in dashboard if row["board_attention"].lower() == "yes"]
    complaint_row = next((r for r in dashboard if "complaint" in r["area"].lower()), None)
    complaints_reduction = None
    if complaint_row:
        current = safe_float(complaint_row["current"])
        previous = safe_float(complaint_row["previous"])
        if current is not None and previous:
            complaints_reduction = (previous - current) / previous

    return {
        "summary": summary,
        "dashboard": dashboard,
        "kpis": kpis,
        "tracker": tracker,
        "counts": counts,
        "may_counts": may_counts,
        "issues": issues,
        "non_compliant": non_compliant,
        "wip": wip,
        "owners": owner_rows,
        "total": total,
        "board_attention": board_attention,
        "complaints_reduction": complaints_reduction,
    }


def autosize_columns(ws):
    for col_idx in range(1, ws.max_column + 1):
        max_len = 10
        for cell in ws.iter_cols(min_col=col_idx, max_col=col_idx):
            for c in cell:
                max_len = max(max_len, len(clean_text(c.value)) + 2)
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len, 48)


def style_table(ws, header_row=1):
    navy = "0D4D2D"
    blue = "1C7A48"
    fill = PatternFill("solid", fgColor=navy)
    thin = Side(style="thin", color="D9E2EF")
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for cell in ws[header_row]:
        cell.fill = fill
        cell.font = Font(color="FFFFFF", bold=True)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.freeze_panes = ws.cell(header_row + 1, 1).coordinate
    ws.sheet_view.showGridLines = False
    return blue


def build_excel_toc(wb, data):
    ws = wb.create_sheet("Table of Contents", 1)
    ws.sheet_view.showGridLines = False
    ws.merge_cells("A1:E1")
    ws["A1"] = "Table of Contents | ARCC Board Compliance Pack Q2 2026"
    ws["A1"].font = Font(bold=True, size=16, color="FFFFFF")
    ws["A1"].fill = PatternFill("solid", fgColor="0D4D2D")
    ws["A1"].alignment = Alignment(horizontal="center")
    ws.append([])
    ws.append(["Section", "Excel Tab", "Purpose", "Board Use", "PPT Slide / Page"])
    rows = [
        ("1", "Board Dashboard", "At-a-glance board dashboard with charts/cards.", "Start here for quick review.", "Slides 3, 9"),
        ("2", "Executive Narrative", "Structured executive verdict, Q2 highlights, policy/governance and next 90 days.", "Supports speaking notes.", "Slides 5-7, 16"),
        ("3", "Board Summary", "Summary metrics and board attention areas.", "Confirm top matters.", "Slide 9"),
        ("4", "Tracker Status", "May vs June status mix.", "Read trend and status split.", "Slide 10"),
        ("5", "KPI Dashboard", "KPI Target/Actual with RAG indicator.", "Monitor performance by indicator.", "Slide 3"),
        ("6", "Open Issues", "Non-compliant and WIP items.", "See remediation universe.", "Slide 14"),
        ("7", "Action Log", "Owner, target date, evidence and escalation fields for open items.", "Track closure after meeting.", "Slides 14, 28-30"),
        ("8", "Owner Heatmap", "Open items by owner.", "Focus accountability.", "Slide 3"),
        ("9", "Detailed Tracker", "Full consolidated tracker.", "Appendix / drill-down.", "Slides 19-27"),
        ("10", "Appendix", "Source documents, assumptions and appendix index.", "Document control.", "Slides 17-30"),
    ]
    for row in rows:
        ws.append(row)
    style_table(ws, header_row=3)
    autosize_columns(ws)


def build_excel_narrative(wb, data):
    ws = wb.create_sheet("Executive Narrative", 2)
    ws.append(["Section", "Board-Ready Content"])
    narrative = [
        ("Executive Verdict", "The Bank maintained a satisfactory compliance posture during Q2 2026. All key regulatory filings highlighted in the executive summary were submitted, and no regulatory sanctions were recorded. The tracker remains satisfactory overall, but 31 open items require dated remediation and evidence-based closure."),
        ("Executive Verdict - Board Attention", "Board attention should focus on AML/CFT implementation maturity, policy implementation, regulatory examination exceptions, technology governance and closure of non-compliant/WIP tracker items."),
        ("Q2 Highlights", "CBN statutory returns, Cybersecurity Return, AML/CFT Implementation Roadmap, Sanctions Designation Report and annual CRS Filing were submitted. AML/CFT Policy was approved, Board Charter received CBN approval, Archival Policy and RSA Policy were developed, and three EFCC requests were concluded."),
        ("Policy & Governance", "AML/CFT Policy and Board Charter are complete/approved. Archival and RSA Policies have been developed. The IT Governance Framework is in progress and will incorporate AI Governance, Cybersecurity Governance, Access Control, BYOD, Third-Party Technology Risk, Information Security and Technology Governance."),
        ("Next 90-Day Focus", "Approve and monitor a dated remediation plan for all 22 non-compliant and 9 WIP items; complete IT Governance Framework; advance AML/CFT Roadmap implementation; close regulatory filing backlogs; and report monthly progress to Management with quarterly escalation to ARCC."),
        ("Conclusion", "The quarter closed with a satisfactory compliance posture and no regulatory sanctions. The required board response is disciplined oversight of open remediation items, evidence-based closure and continued monitoring of emerging technology and AML/CFT expectations."),
    ]
    for row in narrative:
        ws.append(row)
    style_table(ws)
    ws.column_dimensions["A"].width = 28
    ws.column_dimensions["B"].width = 105


def build_excel_appendix(wb, data):
    ws = wb.create_sheet("Appendix")
    ws.append(["Appendix Item", "Description"])
    rows = [
        ("A", "Source: EXECUTIVE SUMMARY.docx"),
        ("B", "Source: BOARD REPORT.xlsx"),
        ("C", "Source: JUNE 2026 CONSOLIDATED TRACKER.xlsx"),
        ("D", f"Full consolidated tracker included in Detailed Tracker tab ({data['total']} items)."),
        ("E", f"Open remediation universe included in Action Log ({len(data['issues'])} open items)."),
        ("F", "PPT appendix contains consolidated tracker detail slides split for readability."),
        ("G", "RAG basis: Green = target met, Amber = close to target/monitoring required, Red = target missed or outstanding issue against nil target."),
    ]
    for row in rows:
        ws.append(row)
    style_table(ws)
    autosize_columns(ws)


def build_excel(data):
    wb = Workbook()
    ws = wb.active
    ws.title = "Board Dashboard"
    ws.sheet_view.showGridLines = False
    ws.merge_cells("A1:I1")
    ws["A1"] = "ARCC Board Compliance Dashboard | Q2 2026"
    ws["A1"].font = Font(bold=True, size=18, color="FFFFFF")
    ws["A1"].fill = PatternFill("solid", fgColor="0D4D2D")
    ws["A1"].alignment = Alignment(horizontal="center")
    ws.merge_cells("A2:I2")
    ws["A2"] = "At-a-glance view for Board oversight: status mix, RAG indicators, open remediation items and key management focus."
    ws["A2"].font = Font(italic=True, color="375F46")
    ws["A2"].alignment = Alignment(horizontal="center")

    dashboard_cards = [
        ("A4:B6", "Overall Status", "Satisfactory", "No regulatory sanctions recorded", "008040"),
        ("C4:D6", "Tracker Items", data["total"], "Consolidated obligations reviewed", "1C7A48"),
        ("E4:F6", "Compliant", data["counts"]["Compliant"], pct(data["counts"]["Compliant"] / data["total"], 0), "008040"),
        ("G4:H6", "Open Items", data["counts"]["Non Compliant"] + data["counts"]["W.I.P"], "Non-compliant + WIP", "D68D1E"),
    ]
    for cell_range, title, value, note, color in dashboard_cards:
        ws.merge_cells(cell_range)
        cell = ws[cell_range.split(":")[0]]
        cell.value = f"{title}\n{value}\n{note}"
        cell.fill = PatternFill("solid", fgColor="F4FAF6")
        cell.font = Font(bold=True, color=color, size=13)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for row in ws[cell_range]:
            for c in row:
                c.border = Border(
                    left=Side(style="thin", color="B7D7C2"),
                    right=Side(style="thin", color="B7D7C2"),
                    top=Side(style="thin", color="B7D7C2"),
                    bottom=Side(style="thin", color="B7D7C2"),
                )

    ws.append([])
    ws["A9"] = "Status"
    ws["B9"] = "Count"
    ws["C9"] = "Share"
    for status in ["Compliant", "Non Compliant", "W.I.P", "N/A"]:
        ws.append([status, data["counts"][status], data["counts"][status] / data["total"]])
    for cell in ws[9]:
        cell.fill = PatternFill("solid", fgColor="0D4D2D")
        cell.font = Font(bold=True, color="FFFFFF")
    for c in ws["C"]:
        if c.row >= 10 and c.row <= 13:
            c.number_format = "0%"

    status_chart = PieChart()
    status_chart.title = "June 2026 Status Mix"
    status_chart.add_data(Reference(ws, min_col=2, min_row=10, max_row=13), titles_from_data=False)
    status_chart.set_categories(Reference(ws, min_col=1, min_row=10, max_row=13))
    status_chart.height = 7
    status_chart.width = 9
    ws.add_chart(status_chart, "E9")

    ws["A17"] = "KPI"
    ws["B17"] = "Target"
    ws["C17"] = "Actual"
    ws["D17"] = "RAG"
    for row in data["kpis"]:
        rag, _ = rag_for_kpi(row["target"], row["actual"], row["status"])
        ws.append([row["kpi"], row["target"], row["actual"], rag])
    for cell in ws[17]:
        cell.fill = PatternFill("solid", fgColor="0D4D2D")
        cell.font = Font(bold=True, color="FFFFFF")
    for r in range(18, 18 + len(data["kpis"])):
        rag = clean_text(ws.cell(r, 4).value)
        fill = "D9EAD3" if rag == "Green" else "FFF2CC" if rag == "Amber" else "F4CCCC"
        ws.cell(r, 4).fill = PatternFill("solid", fgColor=fill)
        ws.cell(r, 4).font = Font(bold=True)

    ws["F17"] = "Owner"
    ws["G17"] = "Open Items"
    for idx, owner in enumerate(data["owners"][:8], 18):
        ws.cell(idx, 6).value = owner["owner"]
        ws.cell(idx, 7).value = owner["open"]
    for c in ("F17", "G17"):
        ws[c].fill = PatternFill("solid", fgColor="0D4D2D")
        ws[c].font = Font(bold=True, color="FFFFFF")
    owner_chart = BarChart()
    owner_chart.type = "bar"
    owner_chart.title = "Top Owners by Open Items"
    owner_chart.add_data(Reference(ws, min_col=7, min_row=17, max_row=25), titles_from_data=True)
    owner_chart.set_categories(Reference(ws, min_col=6, min_row=18, max_row=25))
    owner_chart.height = 7
    owner_chart.width = 9
    ws.add_chart(owner_chart, "E27")
    autosize_columns(ws)
    ws.column_dimensions["A"].width = 34
    ws.column_dimensions["F"].width = 34

    build_excel_toc(wb, data)
    build_excel_narrative(wb, data)

    ws = wb.create_sheet("Board Summary")
    ws.title = "Board Summary"
    ws.append(["ARCC Board Compliance Pack", "Q2 2026"])
    ws.append(["Overall status", "Satisfactory, with defined remediation priorities"])
    ws.append(["Total tracker items", data["total"]])
    ws.append(["Compliant", data["counts"]["Compliant"]])
    ws.append(["Non-compliant", data["counts"]["Non Compliant"]])
    ws.append(["Work in progress", data["counts"]["W.I.P"]])
    ws.append(["N/A", data["counts"]["N/A"]])
    ws.append(["Compliance rate", data["counts"]["Compliant"] / data["total"]])
    ws.append(["Customer complaints reduction", data["complaints_reduction"] or ""])
    ws.append(["No regulatory sanctions recorded", "Yes"])
    ws.append([])
    ws.append(["Board Attention Area", "Current Status", "Comment"])
    for item in data["board_attention"]:
        ws.append([item["area"], format_metric(item["current"]), item["comment"]])
    for row in ws.iter_rows(min_row=1, max_row=10):
        row[0].font = Font(bold=True, color="0D4D2D")
    ws["B8"].number_format = "0%"
    ws["B9"].number_format = "0%"
    style_table(ws, header_row=12)
    autosize_columns(ws)

    ws = wb.create_sheet("Tracker Status")
    ws.append(["Status", "May 2026", "June 2026", "June Share"])
    for status in ["Compliant", "Non Compliant", "W.I.P", "N/A"]:
        ws.append([status, data["may_counts"][status], data["counts"][status], data["counts"][status] / data["total"]])
    style_table(ws)
    for c in ws["D"]:
        if c.row > 1:
            c.number_format = "0%"
    chart = PieChart()
    labels = Reference(ws, min_col=1, min_row=2, max_row=5)
    values = Reference(ws, min_col=3, min_row=2, max_row=5)
    chart.add_data(values, titles_from_data=False)
    chart.set_categories(labels)
    chart.title = "June Status Mix"
    chart.dataLabels = chart.dataLabels or None
    ws.add_chart(chart, "F2")
    autosize_columns(ws)

    ws = wb.create_sheet("KPI Dashboard")
    ws.append(["KPI", "Target", "Actual", "RAG Indicator", "RAG Note"])
    for row in data["kpis"]:
        rag, rag_note = rag_for_kpi(row["target"], row["actual"], row["status"])
        ws.append([row["kpi"], row["target"], row["actual"], rag, rag_note])
    style_table(ws)
    rag_fill = {
        "Green": "D9EAD3",
        "Amber": "FFF2CC",
        "Red": "F4CCCC",
    }
    rag_font = {
        "Green": "0D4D2D",
        "Amber": "7F6000",
        "Red": "990000",
    }
    for row_idx in range(2, ws.max_row + 1):
        rag = clean_text(ws.cell(row_idx, 4).value)
        ws.cell(row_idx, 4).fill = PatternFill("solid", fgColor=rag_fill.get(rag, "FFFFFF"))
        ws.cell(row_idx, 4).font = Font(bold=True, color=rag_font.get(rag, "0D4D2D"))
    kpi_chart = BarChart()
    kpi_chart.title = "KPI Target vs Actual"
    kpi_chart.y_axis.title = "KPI"
    kpi_chart.x_axis.title = "Score / Count"
    kpi_chart.type = "bar"
    kpi_chart.add_data(Reference(ws, min_col=2, max_col=3, min_row=1, max_row=ws.max_row), titles_from_data=True)
    kpi_chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ws.max_row))
    kpi_chart.height = 8
    kpi_chart.width = 12
    ws.add_chart(kpi_chart, "G2")
    autosize_columns(ws)

    ws = wb.create_sheet("Open Issues")
    ws.append(["S/N", "Item", "Owner", "Recipient", "Status", "Timeline", "Remarks", "Proposed Board-Level Action"])
    for row in data["issues"]:
        action = "Management to confirm remediation owner, target date, and escalation requirement."
        if "capital" in row["item"].lower() or "liquidity" in row["item"].lower() or "loan" in row["item"].lower():
            action = "Prioritise prudential remediation and provide monthly progress reporting to ARCC."
        elif "financial statement" in row["item"].lower() or "audited" in row["item"].lower():
            action = "Track dependency on CBN approval and prepare publication/filling immediately on approval."
        elif "report" in row["item"].lower() or "return" in row["item"].lower():
            action = "Close filing backlog and evidence submission in the compliance tracker."
        ws.append([row["sn"], row["item"], row["responsibility"], row["recipient"], row["june_status"], row["timeline"], row["remarks"], action])
    style_table(ws)
    autosize_columns(ws)

    ws = wb.create_sheet("Action Log")
    ws.append([
        "S/N",
        "Open Tracker Item",
        "Owner",
        "Current Status",
        "Target Date",
        "Evidence Required",
        "Escalation Status",
        "Management Update",
        "Board / ARCC Note",
    ])
    for row in data["issues"]:
        item_lower = row["item"].lower()
        evidence = "Evidence of closure, approval, filing, payment, meeting minutes or regulatory acknowledgement."
        target_date = "To be confirmed by owner"
        escalation = "Management"
        board_note = "Track to closure and escalate overdue items to ARCC."
        if "capital" in item_lower or "liquidity" in item_lower or "loan" in item_lower or "obligor" in item_lower:
            evidence = "Prudential ratio computation, remediation plan, management approval and monthly progress evidence."
            escalation = "Board / ARCC"
            board_note = "Requires monthly remediation reporting."
        elif "audited" in item_lower or "financial statement" in item_lower or "afs" in item_lower:
            evidence = "CBN approval evidence, publication/filing proof and regulatory acknowledgement where applicable."
            escalation = "Management / Board"
            board_note = "Track external dependency and readiness to file immediately on approval."
        elif "report" in item_lower or "return" in item_lower or "credit bureau" in item_lower or "str/ctr" in item_lower:
            evidence = "Submission receipt, regulator portal confirmation or acknowledged email."
            escalation = "Compliance / Management"
            board_note = "Close filing backlog and evidence submission."
        elif "committee" in item_lower or "agm" in item_lower:
            evidence = "Approved meeting date, agenda, attendance register and signed minutes."
            escalation = "Management / Board"
            board_note = "Confirm governance calendar and evidence completion."
        ws.append([
            row["sn"],
            row["item"],
            row["responsibility"],
            row["june_status"],
            target_date,
            evidence,
            escalation,
            "",
            board_note,
        ])
    style_table(ws)
    status_fills = {
        "Non Compliant": "F4CCCC",
        "W.I.P": "FFF2CC",
    }
    for row_idx in range(2, ws.max_row + 1):
        status = clean_text(ws.cell(row_idx, 4).value)
        ws.cell(row_idx, 4).fill = PatternFill("solid", fgColor=status_fills.get(status, "FFFFFF"))
        ws.cell(row_idx, 4).font = Font(bold=True, color="990000" if status == "Non Compliant" else "7F6000")
    autosize_columns(ws)

    ws = wb.create_sheet("Owner Heatmap")
    ws.append(["Owner", "Compliant", "Non-Compliant", "W.I.P", "N/A", "Open Items"])
    for row in data["owners"]:
        ws.append([row["owner"], row["compliant"], row["non_compliant"], row["wip"], row["na"], row["open"]])
    style_table(ws)
    chart = BarChart()
    chart.type = "bar"
    chart.title = "Open Items by Owner"
    chart.y_axis.title = "Owner"
    chart.x_axis.title = "Open Items"
    values = Reference(ws, min_col=6, min_row=1, max_row=min(ws.max_row, 12))
    labels = Reference(ws, min_col=1, min_row=2, max_row=min(ws.max_row, 12))
    chart.add_data(values, titles_from_data=True)
    chart.set_categories(labels)
    ws.add_chart(chart, "H2")
    autosize_columns(ws)

    ws = wb.create_sheet("Detailed Tracker")
    ws.append(["S/N", "Item", "Responsibility", "Recipient", "Frequency", "Timeline", "Remarks", "May Status", "June Status"])
    for row in data["tracker"]:
        ws.append([row["sn"], row["item"], row["responsibility"], row["recipient"], row["frequency"], row["timeline"], row["remarks"], row["may_status"], row["june_status"]])
    style_table(ws)
    autosize_columns(ws)

    build_excel_appendix(wb, data)

    wb.save(OUT_XLSX)


def set_run_font(run, size=18, bold=False, color=COLORS["navy"]):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=COLORS["navy"], align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return shape


def add_header(slide, title, kicker="Audit, Risk and Compliance Committee | Q2 2026"):
    add_textbox(slide, Inches(0.55), Inches(0.28), Inches(8.8), Inches(0.25), kicker.upper(), 8, True, COLORS["teal"])
    add_textbox(slide, Inches(0.55), Inches(0.55), Inches(9.3), Inches(0.5), title, 22, True, COLORS["navy"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.08), Inches(9.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()


def add_footer(slide, page):
    add_textbox(slide, Inches(8.8), Inches(7.1), Inches(1.0), Inches(0.18), str(page), 8, False, COLORS["grey"], PP_ALIGN.RIGHT)


def add_bullets(slide, bullets, left, top, width, height, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = COLORS["navy"]
        p.space_after = Pt(8)
    return box


def add_metric_card(slide, left, top, width, height, label, value, note="", color=COLORS["blue"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light"]
    shape.line.color.rgb = RGBColor(218, 225, 235)
    shape.text_frame.clear()
    add_textbox(slide, left + Inches(0.12), top + Inches(0.1), width - Inches(0.24), Inches(0.2), label.upper(), 8, True, COLORS["grey"])
    add_textbox(slide, left + Inches(0.12), top + Inches(0.34), width - Inches(0.24), Inches(0.38), value, 22, True, color)
    if note:
        add_textbox(slide, left + Inches(0.12), top + Inches(0.78), width - Inches(0.24), Inches(0.22), note, 8, False, COLORS["grey"])


def add_table(slide, headers, rows, left, top, width, height, font_size=9):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS["white"]
            p.font.name = "Aptos"
    for r_idx, row in enumerate(rows, 1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = clean_text(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r_idx % 2 else RGBColor(246, 248, 251)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["navy"]
                p.font.name = "Aptos"
    return table


def add_status_chart(slide, data, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = ["Compliant", "Non-Compliant", "WIP", "N/A"]
    chart_data.add_series("June 2026", (
        data["counts"]["Compliant"],
        data["counts"]["Non Compliant"],
        data["counts"]["W.I.P"],
        data["counts"]["N/A"],
    ))
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(9)


def chunked(items, size):
    for start in range(0, len(items), size):
        yield start, items[start:start + size]


def clip(value, max_len=92):
    text = clean_text(value)
    return text if len(text) <= max_len else f"{text[:max_len - 3]}..."


def add_report_toc_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Table of Contents")
    rows = [
        ("1", "Executive Dashboard", "Board Dashboard and KPI RAG", "3"),
        ("2", "Executive Report", "Verdict, Q2 Highlights, Policy & Governance", "5-7"),
        ("3", "Compliance Review", "Compliance dashboard, tracker status and remediation themes", "9-12"),
        ("4", "Management Actions", "Open issues, action log and matters for ARCC", "14-15"),
        ("5", "Forward Focus", "Next 90 days and conclusion", "16-17"),
        ("6", "Appendix", "Detailed consolidated tracker and Excel pack reference", "18-30"),
        ("7", "Closing", "Thank You", "31"),
    ]
    add_table(slide, ["No.", "Section", "Coverage", "Page"], rows, Inches(0.65), Inches(1.35), Inches(8.7), Inches(4.8), 10)
    add_textbox(slide, Inches(0.8), Inches(6.15), Inches(8.4), Inches(0.45), "Excel pack alignment: Table of Contents, Board Dashboard, Executive Narrative, KPI Dashboard, Action Log, Detailed Tracker and Appendix tabs mirror this report structure.", 11, True, COLORS["teal"])
    add_footer(slide, len(prs.slides))


def add_section_divider(prs, blank, title, subtitle):
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(2.25), Inches(8.3), Inches(0.7), title, 30, True, COLORS["white"])
    add_textbox(slide, Inches(0.82), Inches(3.05), Inches(7.6), Inches(0.5), subtitle, 15, False, RGBColor(222, 240, 229))
    add_footer(slide, len(prs.slides))


def build_ppt_v4(data):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    rag_counts = Counter(rag_for_kpi(row["target"], row["actual"], row["status"])[0] for row in data["kpis"])
    open_items = data["counts"]["Non Compliant"] + data["counts"]["W.I.P"]

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(0.75), Inches(8.6), Inches(0.35), "BOARD COMPLIANCE REPORT", 13, True, RGBColor(190, 232, 203))
    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(8.4), Inches(1.2), "Audit, Risk and Compliance Committee", 34, True, COLORS["white"])
    add_textbox(slide, Inches(0.7), Inches(2.6), Inches(7.3), Inches(0.55), "Q2 2026 Compliance Update | Reporting period ended 30 June 2026", 16, False, RGBColor(226, 243, 233))
    add_metric_card(slide, Inches(0.7), Inches(4.2), Inches(2.0), Inches(1.05), "Overall Status", "Satisfactory", "No sanctions recorded", COLORS["green"])
    add_metric_card(slide, Inches(2.95), Inches(4.2), Inches(2.0), Inches(1.05), "Tracker Items", str(data["total"]), "Consolidated obligations", COLORS["teal"])
    add_metric_card(slide, Inches(5.2), Inches(4.2), Inches(2.0), Inches(1.05), "Compliant", str(data["counts"]["Compliant"]), pct(data["counts"]["Compliant"] / data["total"], 0), COLORS["green"])
    add_metric_card(slide, Inches(7.45), Inches(4.2), Inches(2.0), Inches(1.05), "Open Items", str(open_items), "For remediation", COLORS["amber"])

    add_report_toc_slide(prs, blank)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Board Dashboard At A Glance")
    add_metric_card(slide, Inches(0.55), Inches(1.3), Inches(1.75), Inches(1.0), "Compliance Rate", pct(data["counts"]["Compliant"] / data["total"], 0), f"{data['counts']['Compliant']} of {data['total']} items", COLORS["green"])
    add_metric_card(slide, Inches(2.45), Inches(1.3), Inches(1.75), Inches(1.0), "Open Items", str(open_items), "22 red + 9 WIP", COLORS["amber"])
    add_metric_card(slide, Inches(4.35), Inches(1.3), Inches(1.75), Inches(1.0), "KPI RAG", f"{rag_counts['Green']}G / {rag_counts['Amber']}A / {rag_counts['Red']}R", "Indicator sheet updated", COLORS["amber"] if rag_counts["Red"] else COLORS["green"])
    add_metric_card(slide, Inches(6.25), Inches(1.3), Inches(1.75), Inches(1.0), "Complaints", "51", "Down from 218", COLORS["green"])
    add_metric_card(slide, Inches(8.15), Inches(1.3), Inches(1.3), Inches(1.0), "Sanctions", "Nil", "Q2 2026", COLORS["green"])
    add_status_chart(slide, data, Inches(0.65), Inches(2.65), Inches(4.25), Inches(3.1))
    add_table(slide, ["Top Owner", "Open"], [(row["owner"], row["open"]) for row in data["owners"][:5]], Inches(5.2), Inches(2.65), Inches(2.0), Inches(2.05), 9)
    add_table(slide, ["RAG", "Count"], [("Green", rag_counts["Green"]), ("Amber", rag_counts["Amber"]), ("Red", rag_counts["Red"])], Inches(7.45), Inches(2.65), Inches(1.75), Inches(1.5), 10)
    add_bullets(slide, [
        "Board view: satisfactory overall posture with clear remediation accountability.",
        "Excel opens on a dashboard tab, with TOC and Appendix tabs aligned to this deck.",
        "Full consolidated tracker appears in the appendix slides for complete board reporting."
    ], Inches(5.2), Inches(5.05), Inches(4.1), Inches(1.15), 10)
    add_footer(slide, len(prs.slides))

    add_section_divider(prs, blank, "Executive Report", "Board-level narrative and decisions required")

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Executive Verdict")
    add_table(slide, ["Verdict Area", "Board Message"], [
        ("Overall posture", "Satisfactory Q2 2026 compliance posture; no regulatory sanctions recorded."),
        ("Regulatory discipline", "All key regulatory filings highlighted in the executive summary were submitted."),
        ("Tracker interpretation", f"{data['counts']['Compliant']} of {data['total']} consolidated tracker items are compliant; {open_items} items remain open for remediation."),
        ("Board attention", "Focus on AML/CFT implementation maturity, policy implementation, regulatory examination exceptions, technology governance and dated closure of open items."),
    ], Inches(0.7), Inches(1.4), Inches(8.6), Inches(4.3), 11)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Q2 Highlights")
    add_table(slide, ["Theme", "Achievements"], [
        ("Regulatory Returns", "CBN statutory returns, Cybersecurity Return, AML/CFT Roadmap, Sanctions Designation Report and CRS Filing submitted."),
        ("Governance", "AML/CFT Policy approved; Board Charter approved by CBN; Archival Policy and RSA Policy developed."),
        ("Regulatory Engagement", "Three EFCC requests investigated and appropriate responses provided."),
        ("Capacity Building", "Internal and external AML/CFT training completed, plus staff conduct and discipline training."),
    ], Inches(0.7), Inches(1.35), Inches(8.6), Inches(4.4), 11)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Policy & Governance")
    add_table(slide, ["Policy / Framework", "Q2 Status", "Board Relevance"], [
        ("AML/CFT Policy", "Approved", "Strengthens financial crime compliance framework."),
        ("Board Charter", "Approved by CBN", "Confirms governance framework alignment."),
        ("Archival Policy", "Developed", "Supports record retention and evidence management."),
        ("RSA Policy", "Developed", "Supports retirement savings account governance."),
        ("IT Governance Framework", "In progress", "Will incorporate AI, cybersecurity, access control, BYOD, third-party technology risk, information security and technology governance."),
    ], Inches(0.55), Inches(1.3), Inches(8.9), Inches(4.75), 9)
    add_footer(slide, len(prs.slides))

    add_section_divider(prs, blank, "Compliance Review", "Dashboard, tracker status and remediation themes")

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Compliance Dashboard")
    add_metric_card(slide, Inches(0.65), Inches(1.45), Inches(2.0), Inches(1.1), "Reg. Returns", "95%", "Previous: 90%", COLORS["green"])
    add_metric_card(slide, Inches(2.85), Inches(1.45), Inches(2.0), Inches(1.1), "KYC Compliance", "80%", "Previous: 38%", COLORS["green"])
    add_metric_card(slide, Inches(5.05), Inches(1.45), Inches(2.0), Inches(1.1), "Complaints", "51", "Down from 218", COLORS["green"])
    add_metric_card(slide, Inches(7.25), Inches(1.45), Inches(2.0), Inches(1.1), "Exam Exceptions", "8", "Board attention", COLORS["amber"])
    add_table(slide, ["Area", "Status", "Board Attention"], [
        ("AML/CFT Compliance", "Roadmap submitted; vendor solution engagement in progress", "Yes"),
        ("Internal Policy Compliance", "Policies updated/developed; implementation and approval tracking required", "Yes"),
        ("Regulatory Inspection", "CBN target-based examination held", "Yes"),
        ("Regulatory Examination Exceptions", "Most items resolved; board-level items remain", "Yes"),
    ], Inches(0.65), Inches(3.0), Inches(8.7), Inches(2.2), 10)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Consolidated Tracker Status")
    add_status_chart(slide, data, Inches(0.65), Inches(1.45), Inches(4.6), Inches(3.3))
    add_metric_card(slide, Inches(5.55), Inches(1.45), Inches(1.75), Inches(1.0), "Compliant", str(data["counts"]["Compliant"]), pct(data["counts"]["Compliant"] / data["total"], 0), COLORS["green"])
    add_metric_card(slide, Inches(7.55), Inches(1.45), Inches(1.75), Inches(1.0), "Non-Compliant", str(data["counts"]["Non Compliant"]), pct(data["counts"]["Non Compliant"] / data["total"], 0), COLORS["red"])
    add_metric_card(slide, Inches(5.55), Inches(2.75), Inches(1.75), Inches(1.0), "WIP", str(data["counts"]["W.I.P"]), pct(data["counts"]["W.I.P"] / data["total"], 0), COLORS["amber"])
    add_metric_card(slide, Inches(7.55), Inches(2.75), Inches(1.75), Inches(1.0), "N/A", str(data["counts"]["N/A"]), pct(data["counts"]["N/A"] / data["total"], 0), COLORS["grey"])
    add_bullets(slide, [
        "The full consolidated tracker is included in the appendix section of this PPT.",
        "The Excel Detailed Tracker tab contains the same full tracker for filtering and drill-down.",
        "Management should treat all non-compliant and WIP items as the remediation universe."
    ], Inches(5.55), Inches(4.15), Inches(3.8), Inches(1.55), 11)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Priority Remediation Themes")
    add_table(slide, ["Theme", "Examples", "Board-Level Response"], [
        ("Prudential / financial ratios", "Capital, CRR, NPL, single obligor, liquidity and mortgage asset ratios", "Monthly remediation reporting to ARCC"),
        ("Financial/statutory filings", "AFS publication, NITDA levy, NHF remittance, audited accounts", "Track external dependencies and close backlogs"),
        ("Governance cadence", "AGM and selected internal committees", "Confirm dates, owners and evidence of meetings"),
        ("Regulatory and operational returns", "Credit bureau reports, STR/CTR, employee conduct and whistleblowing returns", "Evidence-based filing closure"),
        ("Technology governance", "IT Governance Framework, AI, cybersecurity, access control, BYOD and third-party risk", "Approve framework timetable and monitor implementation"),
    ], Inches(0.6), Inches(1.35), Inches(8.8), Inches(4.6), 9)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "AML/CFT & Financial Crime Compliance")
    add_table(slide, ["Area", "Q2 Position", "Next Management Focus"], [
        ("AML/CFT Roadmap", "Submitted to CBN within required timeline", "Monitor implementation milestones and vendor solution engagement"),
        ("AML/CFT Policy", "Reviewed and approved", "Embed requirements into operating procedures and testing"),
        ("Sanctions Reporting", "NIGSAN and OFAC report submitted on 29 June 2026", "Continue periodic screening evidence and exception reporting"),
        ("Training", "Internal and external AML/CFT training completed", "Maintain coverage records and refresh schedule"),
        ("EFCC Requests", "Three requests investigated and responded to", "Maintain response evidence and lessons learned"),
    ], Inches(0.6), Inches(1.35), Inches(8.8), Inches(4.6), 10)
    add_footer(slide, len(prs.slides))

    add_section_divider(prs, blank, "Management Actions", "Open issues, action log and matters for ARCC")

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Management Action Tracker")
    rows = [(clip(i["item"], 82), i["responsibility"], i["june_status"]) for i in data["issues"][:8]]
    add_table(slide, ["Open Item", "Owner", "Status"], rows, Inches(0.6), Inches(1.25), Inches(8.8), Inches(4.7), 8)
    add_textbox(slide, Inches(0.75), Inches(6.25), Inches(8.5), Inches(0.42), "Excel Action Log contains all open tracker items with owner, target date, evidence required, escalation status and management update fields.", 12, True, COLORS["amber"])
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Matters For ARCC")
    add_table(slide, ["For Noting", "For Discussion", "For Approval"], [
        ("All key Q2 regulatory returns and filings submitted; no sanctions recorded.", "Progress on AML/CFT Roadmap implementation and vendor solution engagement.", "None in the source pack unless newly developed policies are ready for recommendation."),
        ("AML/CFT Policy approved; Board Charter approved by CBN.", "Progress on IT Governance Framework and emerging AI/cybersecurity governance requirements.", ""),
        ("Three EFCC investigations concluded; internal/external AML training completed.", "Closure plan for tracker non-compliant and WIP items, including prudential and filing-related items.", ""),
    ], Inches(0.55), Inches(1.35), Inches(8.9), Inches(4.9), 10)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Focus For The Next 90 Days")
    add_table(slide, ["Priority", "Expected Management Action", "Board Oversight Lens"], [
        ("Open tracker remediation", "Approve and monitor dated remediation plan for all 22 non-compliant and 9 WIP items.", "Owner, target date, evidence and escalation status."),
        ("IT Governance Framework", "Complete framework covering AI, cybersecurity, access control, BYOD and third-party technology risk.", "Timetable, policy approval and implementation evidence."),
        ("AML/CFT Roadmap", "Advance implementation, vendor engagement and milestone evidence.", "Regulatory expectation and control maturity."),
        ("Filing backlogs", "Close regulatory filing backlogs and evidence all submissions.", "Submission proof and exception monitoring."),
        ("Reporting cadence", "Report monthly progress to Management and quarterly progress to ARCC.", "Consistent board visibility."),
    ], Inches(0.6), Inches(1.35), Inches(8.8), Inches(4.9), 9)
    add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Conclusion")
    add_bullets(slide, [
        "Q2 2026 closed with a satisfactory compliance posture and no recorded regulatory sanctions.",
        "The Bank completed important regulatory submissions, policy actions, training and engagement activities.",
        "The key board requirement is disciplined follow-through on the 31 open tracker items through the Action Log.",
        "Management should evidence closure, escalate overdue items and report progress through the agreed governance cadence."
    ], Inches(0.8), Inches(1.45), Inches(8.4), Inches(4.8), 16)
    add_footer(slide, len(prs.slides))

    add_section_divider(prs, blank, "Appendix", "Full consolidated tracker and Excel pack reference")

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Appendix Index")
    add_table(slide, ["Appendix", "Content", "Excel Reference"], [
        ("A", "Consolidated tracker detail", "Detailed Tracker tab"),
        ("B", "Open-item action log", "Action Log tab"),
        ("C", "KPI RAG basis", "KPI Dashboard tab"),
        ("D", "Source document list", "Appendix tab"),
    ], Inches(0.9), Inches(1.55), Inches(8.1), Inches(3.1), 12)
    add_textbox(slide, Inches(0.9), Inches(5.2), Inches(8.1), Inches(0.6), "The appendix slides below include the full consolidated tracker so the PPT report is complete. The Excel pack remains the working document for filtering, updates and management follow-up.", 12, True, COLORS["teal"])
    add_footer(slide, len(prs.slides))

    tracker_rows = [
        (row["sn"], clip(row["item"], 48), clip(row["responsibility"], 22), clip(row["recipient"], 22), row["june_status"])
        for row in data["tracker"]
    ]
    for start, rows in chunked(tracker_rows, 14):
        slide = prs.slides.add_slide(blank)
        add_header(slide, f"Appendix A: Consolidated Tracker ({start + 1}-{start + len(rows)} of {len(tracker_rows)})")
        add_table(slide, ["S/N", "Item", "Owner", "Recipient", "June Status"], rows, Inches(0.35), Inches(1.25), Inches(9.3), Inches(5.75), 6)
        add_footer(slide, len(prs.slides))

    action_rows = [
        (row["sn"], clip(row["item"], 50), clip(row["responsibility"], 22), row["june_status"], "To be confirmed")
        for row in data["issues"]
    ]
    for start, rows in chunked(action_rows, 12):
        slide = prs.slides.add_slide(blank)
        add_header(slide, f"Appendix B: Action Log ({start + 1}-{start + len(rows)} of {len(action_rows)})")
        add_table(slide, ["S/N", "Open Item", "Owner", "Status", "Target Date"], rows, Inches(0.35), Inches(1.25), Inches(9.3), Inches(5.75), 6)
        add_footer(slide, len(prs.slides))

    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(2.55), Inches(8.4), Inches(0.9), "Thank You", 40, True, COLORS["white"], PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.15), Inches(3.5), Inches(7.7), Inches(0.55), "Questions and board guidance", 18, False, RGBColor(226, 243, 233), PP_ALIGN.CENTER)
    add_footer(slide, len(prs.slides))

    prs.save(OUT_PPTX)


def build_ppt(data):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["navy"]
    bg.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(0.75), Inches(8.6), Inches(0.35), "BOARD COMPLIANCE REPORT", 13, True, RGBColor(139, 211, 212))
    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(8.4), Inches(1.2), "Audit, Risk and Compliance Committee", 34, True, COLORS["white"])
    add_textbox(slide, Inches(0.7), Inches(2.6), Inches(7.3), Inches(0.55), "Q2 2026 Compliance Update | Reporting period ended 30 June 2026", 16, False, RGBColor(226, 233, 243))
    add_metric_card(slide, Inches(0.7), Inches(4.2), Inches(2.0), Inches(1.05), "Overall Status", "Satisfactory", "No sanctions recorded", RGBColor(15, 160, 112))
    add_metric_card(slide, Inches(2.95), Inches(4.2), Inches(2.0), Inches(1.05), "Tracker Items", str(data["total"]), "Consolidated obligations", RGBColor(139, 211, 212))
    add_metric_card(slide, Inches(5.2), Inches(4.2), Inches(2.0), Inches(1.05), "Compliant", str(data["counts"]["Compliant"]), pct(data["counts"]["Compliant"] / data["total"], 0), RGBColor(15, 160, 112))
    add_metric_card(slide, Inches(7.45), Inches(4.2), Inches(2.0), Inches(1.05), "Open Items", str(data["counts"]["Non Compliant"] + data["counts"]["W.I.P"]), "For remediation", RGBColor(214, 141, 30))

    # 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Board Dashboard At A Glance")
    open_items = data["counts"]["Non Compliant"] + data["counts"]["W.I.P"]
    rag_counts = Counter(rag_for_kpi(row["target"], row["actual"], row["status"])[0] for row in data["kpis"])
    add_metric_card(slide, Inches(0.55), Inches(1.3), Inches(1.75), Inches(1.0), "Compliance Rate", pct(data["counts"]["Compliant"] / data["total"], 0), f"{data['counts']['Compliant']} of {data['total']} items", COLORS["green"])
    add_metric_card(slide, Inches(2.45), Inches(1.3), Inches(1.75), Inches(1.0), "Open Items", str(open_items), "22 red + 9 WIP", COLORS["amber"])
    add_metric_card(slide, Inches(4.35), Inches(1.3), Inches(1.75), Inches(1.0), "KPI RAG", f"{rag_counts['Green']}G / {rag_counts['Amber']}A / {rag_counts['Red']}R", "Indicator sheet updated", COLORS["green"] if rag_counts["Red"] == 0 else COLORS["amber"])
    add_metric_card(slide, Inches(6.25), Inches(1.3), Inches(1.75), Inches(1.0), "Complaints", "51", "Down from 218", COLORS["green"])
    add_metric_card(slide, Inches(8.15), Inches(1.3), Inches(1.3), Inches(1.0), "Sanctions", "Nil", "Q2 2026", COLORS["green"])
    add_status_chart(slide, data, Inches(0.65), Inches(2.65), Inches(4.25), Inches(3.1))
    add_table(slide, ["Top Owner", "Open"], [(row["owner"], row["open"]) for row in data["owners"][:5]], Inches(5.2), Inches(2.65), Inches(2.0), Inches(2.05), 9)
    add_table(slide, ["RAG", "Count"], [("Green", rag_counts["Green"]), ("Amber", rag_counts["Amber"]), ("Red", rag_counts["Red"])], Inches(7.45), Inches(2.65), Inches(1.75), Inches(1.5), 10)
    add_bullets(slide, [
        "Board view: satisfactory overall posture, with remediation focus on non-compliant and WIP tracker items.",
        "KPI sheet now uses RAG indicators for target-versus-actual interpretation.",
        "Excel pack opens on a dashboard tab for immediate review before drill-down."
    ], Inches(5.2), Inches(5.05), Inches(4.1), Inches(1.15), 10)
    add_footer(slide, len(prs.slides))

    slides = []
    # 3
    slides.append(("Executive Verdict", [
        "The Bank maintained a satisfactory compliance posture during Q2 2026, with all key regulatory filings highlighted in the executive summary submitted.",
        "No regulatory sanctions were recorded during the quarter.",
        "The consolidated tracker shows a material compliance base, but 31 items require management remediation or close monitoring.",
        "Board attention is required on AML/CFT implementation maturity, policy implementation, regulatory examination exceptions, and technology governance."
    ]))
    # 4
    slides.append(("Q2 Highlights", [
        "Submitted CBN statutory returns, Cybersecurity Return, AML/CFT Implementation Roadmap, Sanctions Designation Report, and annual CRS Filing.",
        "AML/CFT Policy reviewed and approved; Board Charter approved by CBN.",
        "Archival Policy and RSA Policy developed for governance strengthening.",
        "Three EFCC requests investigated and responses provided.",
        "Internal and external AML/CFT training completed, alongside staff conduct and discipline training."
    ]))
    for title, bullets in slides:
        slide = prs.slides.add_slide(blank)
        add_header(slide, title)
        add_bullets(slide, bullets, Inches(0.75), Inches(1.45), Inches(8.6), Inches(4.8), 16)
        add_footer(slide, len(prs.slides))

    # 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Compliance Dashboard")
    add_metric_card(slide, Inches(0.65), Inches(1.45), Inches(2.0), Inches(1.1), "Reg. Returns", "95%", "Previous: 90%", COLORS["green"])
    add_metric_card(slide, Inches(2.85), Inches(1.45), Inches(2.0), Inches(1.1), "KYC Compliance", "80%", "Previous: 38%", COLORS["green"])
    add_metric_card(slide, Inches(5.05), Inches(1.45), Inches(2.0), Inches(1.1), "Complaints", "51", "Down from 218", COLORS["green"])
    add_metric_card(slide, Inches(7.25), Inches(1.45), Inches(2.0), Inches(1.1), "Exam Exceptions", "8", "Board attention", COLORS["amber"])
    add_table(slide, ["Area", "Status", "Board Attention"], [
        ("AML/CFT Compliance", "Roadmap submitted; vendor solution engagement in progress", "Yes"),
        ("Internal Policy Compliance", "Policies updated/developed; implementation and approval tracking required", "Yes"),
        ("Regulatory Inspection", "CBN target-based examination held", "Yes"),
        ("Regulatory Examination Exceptions", "Most items resolved; board-level items remain", "Yes"),
    ], Inches(0.65), Inches(3.0), Inches(8.7), Inches(2.2), 10)
    add_footer(slide, len(prs.slides))

    # 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Consolidated Tracker Status")
    add_status_chart(slide, data, Inches(0.65), Inches(1.45), Inches(4.6), Inches(3.3))
    add_metric_card(slide, Inches(5.55), Inches(1.45), Inches(1.75), Inches(1.0), "Compliant", str(data["counts"]["Compliant"]), pct(data["counts"]["Compliant"] / data["total"], 0), COLORS["green"])
    add_metric_card(slide, Inches(7.55), Inches(1.45), Inches(1.75), Inches(1.0), "Non-Compliant", str(data["counts"]["Non Compliant"]), pct(data["counts"]["Non Compliant"] / data["total"], 0), COLORS["red"])
    add_metric_card(slide, Inches(5.55), Inches(2.75), Inches(1.75), Inches(1.0), "WIP", str(data["counts"]["W.I.P"]), pct(data["counts"]["W.I.P"] / data["total"], 0), COLORS["amber"])
    add_metric_card(slide, Inches(7.55), Inches(2.75), Inches(1.75), Inches(1.0), "N/A", str(data["counts"]["N/A"]), pct(data["counts"]["N/A"] / data["total"], 0), COLORS["grey"])
    add_bullets(slide, [
        "Management should treat the 31 open items as the immediate remediation universe.",
        "Priority owners by open items: FINCON/Treasury, Credit & Risk Management, Compliance, General Service, and Operations.",
        "The tracker should move from status reporting to dated action ownership for all non-compliant and WIP items."
    ], Inches(5.55), Inches(4.15), Inches(3.8), Inches(1.55), 11)
    add_footer(slide, len(prs.slides))

    # 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Priority Remediation Themes")
    add_table(slide, ["Theme", "Examples", "Board-Level Response"], [
        ("Prudential / financial ratios", "Capital, CRR, NPL, single obligor, liquidity and mortgage asset ratios", "Monthly remediation reporting to ARCC"),
        ("Financial/statutory filings", "AFS publication, NITDA levy, NHF remittance, audited accounts", "Track external dependencies and close backlogs"),
        ("Governance cadence", "AGM and selected internal committees", "Confirm dates, owners and evidence of meetings"),
        ("Regulatory and operational returns", "Credit bureau reports, STR/CTR, employee conduct and whistleblowing returns", "Evidence-based filing closure"),
        ("Technology governance", "IT Governance Framework, AI, cybersecurity, access control, BYOD and third-party risk", "Approve framework timetable and monitor implementation"),
    ], Inches(0.6), Inches(1.35), Inches(8.8), Inches(4.6), 9)
    add_footer(slide, len(prs.slides))

    # 7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Policy & Governance")
    add_bullets(slide, [
        "AML/CFT Policy approved.",
        "Board Charter approved by CBN.",
        "Archival Policy and RSA Policy developed.",
        "IT Governance Framework in progress and expanded to cover AI Governance, Cybersecurity Governance, Access Control, BYOD, Third-Party Technology Risk, Information Security and Technology Governance."
    ], Inches(0.75), Inches(1.4), Inches(8.7), Inches(4.2), 16)
    add_footer(slide, len(prs.slides))

    # 8
    slide = prs.slides.add_slide(blank)
    add_header(slide, "AML/CFT & Financial Crime Compliance")
    add_table(slide, ["Area", "Q2 Position", "Next Management Focus"], [
        ("AML/CFT Roadmap", "Submitted to CBN within required timeline", "Monitor implementation milestones and vendor solution engagement"),
        ("AML/CFT Policy", "Reviewed and approved", "Embed requirements into operating procedures and testing"),
        ("Sanctions Reporting", "NIGSAN and OFAC report submitted on 29 June 2026", "Continue periodic screening evidence and exception reporting"),
        ("Training", "Internal and external AML/CFT training completed", "Maintain coverage records and refresh schedule"),
        ("EFCC Requests", "Three requests investigated and responded to", "Maintain response evidence and lessons learned"),
    ], Inches(0.6), Inches(1.35), Inches(8.8), Inches(4.6), 10)
    add_footer(slide, len(prs.slides))

    # 9
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Management Action Tracker")
    top_issues = data["issues"][:8]
    rows = [(i["item"], i["responsibility"], i["june_status"]) for i in top_issues]
    add_table(slide, ["Open Item", "Owner", "Status"], rows, Inches(0.6), Inches(1.25), Inches(8.8), Inches(4.7), 8)
    add_textbox(slide, Inches(0.75), Inches(6.25), Inches(8.5), Inches(0.42), "Recommendation implemented in Excel: every open item now has an Action Log line for owner, target date, evidence required and escalation status.", 12, True, COLORS["amber"])
    add_footer(slide, len(prs.slides))

    # 10
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Matters For ARCC")
    add_table(slide, ["For Noting", "For Discussion", "For Approval"], [
        ("All key Q2 regulatory returns and filings submitted; no sanctions recorded.", "Progress on AML/CFT Roadmap implementation and vendor solution engagement.", "None in the source pack unless newly developed policies are ready for recommendation."),
        ("AML/CFT Policy approved; Board Charter approved by CBN.", "Progress on IT Governance Framework and emerging AI/cybersecurity governance requirements.", ""),
        ("Three EFCC investigations concluded; internal/external AML training completed.", "Closure plan for tracker non-compliant and WIP items, including prudential and filing-related items.", ""),
    ], Inches(0.55), Inches(1.35), Inches(8.9), Inches(4.9), 10)
    add_footer(slide, len(prs.slides))

    # 11
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Next 90-Day Focus")
    add_bullets(slide, [
        "Approve and monitor a dated remediation plan for all 22 non-compliant and 9 WIP tracker items.",
        "Complete IT Governance Framework and align it with AI, cybersecurity, access control, BYOD and third-party technology risk expectations.",
        "Advance AML/CFT Roadmap implementation, including vendor engagement and implementation evidence.",
        "Close regulatory filing backlogs and evidence all submissions in the consolidated tracker.",
        "Report monthly progress to Management and quarterly progress to ARCC."
    ], Inches(0.75), Inches(1.4), Inches(8.6), Inches(4.9), 16)
    add_footer(slide, len(prs.slides))

    prs.save(OUT_PPTX)


def pdf_slide(doc, title, bullets=None, table_data=None, metrics=None):
    styles = getSampleStyleSheet()
    h = ParagraphStyle("SlideTitle", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=24, leading=30, textColor=colors.HexColor("#0D4D2D"), alignment=TA_LEFT)
    b = ParagraphStyle("Bullet", parent=styles["BodyText"], fontName="Helvetica", fontSize=13, leading=18, leftIndent=12, bulletIndent=0, textColor=colors.HexColor("#0D4D2D"))
    normal = ParagraphStyle("Normal2", parent=styles["BodyText"], fontName="Helvetica", fontSize=11, leading=15, textColor=colors.HexColor("#0D4D2D"))
    story = [Paragraph(title, h), Spacer(1, 0.18 * inch)]
    if metrics:
        tbl = Table([[Paragraph(k, normal), Paragraph(str(v), normal)] for k, v in metrics], colWidths=[2.2 * inch, 1.6 * inch])
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F5F7FA")),
            ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#D9E2EF")),
            ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#D9E2EF")),
            ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.extend([tbl, Spacer(1, 0.2 * inch)])
    if bullets:
        for item in bullets:
            story.append(Paragraph(f"- {item}", b))
            story.append(Spacer(1, 0.07 * inch))
    if table_data:
        tbl = Table(table_data, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0D4D2D")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("BOX", (0, 0), (-1, -1), 0.25, colors.HexColor("#D9E2EF")),
            ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#D9E2EF")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F6F8FB")]),
        ]))
        story.append(tbl)
    story.append(PageBreak())
    doc.extend(story)


def build_pdf(data):
    story = []
    doc = SimpleDocTemplate(str(OUT_PDF), pagesize=landscape(letter), rightMargin=0.55 * inch, leftMargin=0.55 * inch, topMargin=0.5 * inch, bottomMargin=0.45 * inch)
    rag_counts = Counter(rag_for_kpi(row["target"], row["actual"], row["status"])[0] for row in data["kpis"])
    pdf_slide(story, "ARCC Board Compliance Report | Q2 2026", metrics=[
        ("Overall Status", "Satisfactory"),
        ("Total Tracker Items", data["total"]),
        ("Compliant", f"{data['counts']['Compliant']} ({pct(data['counts']['Compliant'] / data['total'], 0)})"),
        ("Open Items", data["counts"]["Non Compliant"] + data["counts"]["W.I.P"]),
    ], bullets=[
        "Reporting period ended 30 June 2026.",
        "No regulatory sanctions were recorded during the quarter.",
        "This PDF mirrors the board presentation and can be circulated as a meeting pack."
    ])
    pdf_slide(story, "Table of Contents", table_data=[
        ["No.", "Section", "Coverage", "Page"],
        ["1", "Executive Dashboard", "Board Dashboard and KPI RAG", "3"],
        ["2", "Executive Report", "Verdict, Q2 Highlights, Policy & Governance", "5-7"],
        ["3", "Compliance Review", "Compliance dashboard, tracker status and remediation themes", "9-12"],
        ["4", "Management Actions", "Open issues, action log and matters for ARCC", "14-15"],
        ["5", "Forward Focus", "Next 90 days and conclusion", "16-17"],
        ["6", "Appendix", "Consolidated tracker and Excel pack reference", "18-30"],
        ["7", "Closing", "Thank You", "31"],
    ], bullets=[
        "Excel pack alignment: TOC, Board Dashboard, Executive Narrative, KPI Dashboard, Action Log, Detailed Tracker and Appendix tabs mirror this report structure."
    ])
    pdf_slide(story, "Board Dashboard At A Glance", metrics=[
        ("Compliance Rate", pct(data["counts"]["Compliant"] / data["total"], 0)),
        ("Open Items", data["counts"]["Non Compliant"] + data["counts"]["W.I.P"]),
        ("KPI RAG", f"{rag_counts['Green']} Green / {rag_counts['Amber']} Amber / {rag_counts['Red']} Red"),
        ("Customer Complaints", "51, down from 218"),
    ], bullets=[
        "Excel opens on a dashboard tab with status mix, KPI RAG, owner exposure and open remediation items.",
        "KPI sheet now includes a formal RAG Indicator and RAG Note for target-versus-actual interpretation.",
        "PowerPoint includes a front dashboard slide for immediate board-level scanning."
    ])
    pdf_slide(story, "Executive Verdict", bullets=[
        "The Bank maintained a satisfactory compliance posture during Q2 2026.",
        "The consolidated tracker shows 73 compliant items, 22 non-compliant items, 9 WIP items and 5 N/A items.",
        "Board attention is required on AML/CFT implementation maturity, policy implementation, regulatory examination exceptions and technology governance.",
    ])
    pdf_slide(story, "Q2 Highlights", bullets=[
        "Submitted CBN statutory returns, Cybersecurity Return, AML/CFT Roadmap, Sanctions Designation Report and CRS Filing.",
        "AML/CFT Policy approved; Board Charter approved by CBN.",
        "Archival Policy and RSA Policy developed.",
        "Three EFCC requests investigated and responded to.",
        "Internal and external AML/CFT training completed."
    ])
    pdf_slide(story, "Compliance Dashboard", metrics=[
        ("Regulatory Returns", "95% (Previous: 90%)"),
        ("KYC Compliance", "80% (Previous: 38%)"),
        ("Customer Complaints", "51 (Down from 218)"),
        ("Regulatory Exam Exceptions", "8"),
    ], table_data=[
        ["Area", "Status", "Board Attention"],
        ["AML/CFT Compliance", "Roadmap submitted; vendor solution engagement in progress", "Yes"],
        ["Internal Policy Compliance", "Policies updated/developed; implementation and approval tracking required", "Yes"],
        ["Regulatory Inspection", "CBN target-based examination held", "Yes"],
        ["Regulatory Examination Exceptions", "Most items resolved; board-level items remain", "Yes"],
    ])
    pdf_slide(story, "Consolidated Tracker Status", metrics=[
        ("Compliant", data["counts"]["Compliant"]),
        ("Non-Compliant", data["counts"]["Non Compliant"]),
        ("WIP", data["counts"]["W.I.P"]),
        ("N/A", data["counts"]["N/A"]),
    ], bullets=[
        "Management should treat the 31 non-compliant/WIP items as the immediate remediation universe.",
        "Priority owners by open items include FINCON/Treasury, Credit & Risk Management, Compliance, General Service and Operations."
    ])
    pdf_slide(story, "Priority Remediation Themes", table_data=[
        ["Theme", "Examples", "Board-Level Response"],
        ["Prudential / financial ratios", "Capital, CRR, NPL, liquidity and mortgage asset ratios", "Monthly remediation reporting to ARCC"],
        ["Financial/statutory filings", "AFS publication, NITDA levy, NHF remittance, audited accounts", "Track dependencies and close backlogs"],
        ["Governance cadence", "AGM and selected internal committees", "Confirm dates, owners and evidence"],
        ["Regulatory returns", "Credit bureau reports, STR/CTR, employee conduct and whistleblowing returns", "Evidence-based filing closure"],
        ["Technology governance", "IT Governance Framework, AI, cybersecurity, access control and third-party risk", "Approve timetable and monitor implementation"],
    ])
    pdf_slide(story, "Policy & Governance", bullets=[
        "AML/CFT Policy approved.",
        "Board Charter approved by CBN.",
        "Archival Policy and RSA Policy developed.",
        "IT Governance Framework is in progress and is being expanded to cover AI Governance, Cybersecurity Governance, Access Control, BYOD, Third-Party Technology Risk, Information Security and Technology Governance.",
    ])
    pdf_slide(story, "AML/CFT & Financial Crime Compliance", table_data=[
        ["Area", "Q2 Position", "Next Management Focus"],
        ["AML/CFT Roadmap", "Submitted to CBN within required timeline", "Monitor implementation milestones and vendor solution engagement"],
        ["AML/CFT Policy", "Reviewed and approved", "Embed requirements into operating procedures and testing"],
        ["Sanctions Reporting", "NIGSAN and OFAC report submitted on 29 June 2026", "Continue periodic screening evidence and exception reporting"],
        ["Training", "Internal and external AML/CFT training completed", "Maintain coverage records and refresh schedule"],
        ["EFCC Requests", "Three requests investigated and responded to", "Maintain response evidence and lessons learned"],
    ])
    pdf_slide(story, "Management Action Tracker", table_data=[
        ["Open Item", "Owner", "Status"],
        *[[row["item"], row["responsibility"], row["june_status"]] for row in data["issues"][:8]],
    ], bullets=[
        "Recommendation: convert every open tracker item into an action log with owner, target date, evidence required and escalation status."
    ])
    pdf_slide(story, "Matters For ARCC", table_data=[
        ["For Noting", "For Discussion", "For Approval"],
        ["All key Q2 regulatory returns and filings submitted; no sanctions recorded.", "Progress on AML/CFT Roadmap and vendor engagement.", "None in source pack unless policies are ready for recommendation."],
        ["AML/CFT Policy approved; Board Charter approved by CBN.", "Progress on IT Governance Framework and AI/cybersecurity governance.", ""],
        ["Three EFCC investigations concluded; AML training completed.", "Closure plan for non-compliant and WIP tracker items.", ""],
    ])
    pdf_slide(story, "Next 90-Day Focus", bullets=[
        "Approve and monitor a dated remediation plan for all 22 non-compliant and 9 WIP tracker items.",
        "Complete IT Governance Framework and align it with AI, cybersecurity, access control, BYOD and third-party technology risk expectations.",
        "Advance AML/CFT Roadmap implementation and retain evidence of vendor engagement and milestones.",
        "Close regulatory filing backlogs and evidence submissions in the consolidated tracker.",
        "Report monthly progress to Management and quarterly progress to ARCC."
    ])
    pdf_slide(story, "Conclusion", bullets=[
        "Q2 2026 closed with a satisfactory compliance posture and no recorded regulatory sanctions.",
        "The Bank completed important regulatory submissions, policy actions, training and engagement activities.",
        "The key board requirement is disciplined follow-through on the 31 open tracker items through the Action Log.",
        "Management should evidence closure, escalate overdue items and report progress through the agreed governance cadence."
    ])
    pdf_slide(story, "Appendix Index", table_data=[
        ["Appendix", "Content", "Excel Reference"],
        ["A", "Consolidated tracker detail", "Detailed Tracker tab"],
        ["B", "Open-item action log", "Action Log tab"],
        ["C", "KPI RAG basis", "KPI Dashboard tab"],
        ["D", "Source document list", "Appendix tab"],
    ])
    tracker_rows = [
        [row["sn"], clip(row["item"], 58), clip(row["responsibility"], 24), row["june_status"]]
        for row in data["tracker"]
    ]
    for start, rows in chunked(tracker_rows, 20):
        pdf_slide(story, f"Appendix A: Consolidated Tracker ({start + 1}-{start + len(rows)} of {len(tracker_rows)})", table_data=[
            ["S/N", "Item", "Owner", "June Status"],
            *rows,
        ])
    action_rows = [
        [row["sn"], clip(row["item"], 58), clip(row["responsibility"], 24), row["june_status"], "To be confirmed"]
        for row in data["issues"]
    ]
    for start, rows in chunked(action_rows, 18):
        pdf_slide(story, f"Appendix B: Action Log ({start + 1}-{start + len(rows)} of {len(action_rows)})", table_data=[
            ["S/N", "Open Item", "Owner", "Status", "Target Date"],
            *rows,
        ])
    pdf_slide(story, "Thank You", bullets=[
        "Questions and board guidance."
    ])
    if story and isinstance(story[-1], PageBreak):
        story.pop()
    doc.build(story)


def main():
    data = derive_pack_data()
    build_excel(data)
    build_ppt_v4(data)
    build_pdf(data)
    print(f"Created: {OUT_PPTX}")
    print(f"Created: {OUT_XLSX}")
    print(f"Created: {OUT_PDF}")


if __name__ == "__main__":
    main()
